import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import io
import os

from pptx import Presentation
from pptx.util import Inches, Pt

st.set_page_config(page_title="Consolidada IMB", layout="wide")

# =========================
# ESTILO
# =========================
st.markdown("""
<style>
html, body, [class*="css"] {
    font-family: 'Noto Sans', sans-serif;
}

.titulo {
    color:#9F2241;
    font-weight:800;
    text-align:center;
}

thead tr th {
    background-color:#235B4E !important;
    color:white !important;
    font-weight:bold !important;
}

[data-testid="stDataFrame"] {
    border:2px solid #235B4E;
}
</style>
""", unsafe_allow_html=True)

st.markdown(
    "<h1 class='titulo'>📊 Dashboard IMSS BIENESTAR - Abasto</h1>",
    unsafe_allow_html=True
)

# =========================
# ARCHIVO
# =========================
archivo = r"C:\Users\guillermo.ortega\IMSS-BIENESTAR\Emisión y Proveeduría - Consolidada 2026\Consolidada 2026 IMB.xlsb"

@st.cache_data(show_spinner=False)
def cargar():
    df = pd.read_excel(
        archivo,
        engine="pyxlsb",
        usecols=[0, 3, 17, 24, 27, 55]
    )
    df.columns = df.columns.str.strip()
    return df

df = cargar()

# =========================
# COLUMNAS
# =========================
col_orden = df.columns[0]
col_estado = df.columns[1]
col_clave = df.columns[2]
col_precio = df.columns[3]
col_emitidas = df.columns[4]
col_entregadas = df.columns[5]

# =========================
# FORMATOS
# =========================
def fmt(x):
    return f"{int(x):,}"

def fmt_money(x):
    return f"${int(x):,}"

# =========================
# MÉTRICAS
# =========================
def calcular_metricas(df_estado):

    df_estado = df_estado.copy()

    df_estado["Emitido"] = pd.to_numeric(
        df_estado[col_emitidas],
        errors="coerce"
    ).fillna(0)

    df_estado["Entregado"] = pd.to_numeric(
        df_estado[col_entregadas],
        errors="coerce"
    ).fillna(0)

    df_estado[col_precio] = pd.to_numeric(
        df_estado[col_precio],
        errors="coerce"
    ).fillna(0)

    df_estado["Transito"] = (
        df_estado["Emitido"] - df_estado["Entregado"]
    ).clip(lower=0)

    df_estado["Monto_Emitido"] = df_estado["Emitido"] * df_estado[col_precio]
    df_estado["Monto_Entregado"] = df_estado["Entregado"] * df_estado[col_precio]
    df_estado["Monto_Transito"] = df_estado["Monto_Emitido"] - df_estado["Monto_Entregado"]

    p = df_estado["Emitido"].sum()
    e = df_estado["Entregado"].sum()
    t = df_estado["Transito"].sum()

    m_p = df_estado["Monto_Emitido"].sum()
    m_e = df_estado["Monto_Entregado"].sum()
    m_t = df_estado["Monto_Transito"].sum()

    o_total = df_estado[col_orden].nunique()
    c_total = df_estado[col_clave].nunique()

    o_ent = df_estado[df_estado["Entregado"] >= df_estado["Emitido"]][col_orden].nunique()
    o_tran = df_estado[df_estado["Entregado"] < df_estado["Emitido"]][col_orden].nunique()

    c_ent = df_estado[df_estado["Entregado"] >= df_estado["Emitido"]][col_clave].nunique()
    c_tran = df_estado[df_estado["Entregado"] < df_estado["Emitido"]][col_clave].nunique()

    tabla = pd.DataFrame([
        ["Piezas", fmt(p), fmt(t), fmt(e)],
        ["Claves", fmt(c_total), fmt(c_tran), fmt(c_ent)],
        ["Órdenes", fmt(o_total), fmt(o_tran), fmt(o_ent)],
        ["Montos", fmt_money(m_p), fmt_money(m_t), fmt_money(m_e)]
    ], columns=["Métrica", "Emitido", "En tránsito", "Entregado"])

    return {
        "p": p,
        "e": e,
        "t": t,
        "m_p": m_p,
        "m_e": m_e,
        "m_t": m_t,
        "o_total": o_total,
        "o_ent": o_ent,
        "o_tran": o_tran,
        "c_total": c_total,
        "c_ent": c_ent,
        "c_tran": c_tran,
        "tabla": tabla
    }

# =========================
# GRÁFICAS
# =========================
def grafica(titulo, valores, labels, colores):

    fig, ax = plt.subplots(figsize=(10.5, 3.6))

    ax.barh(labels, valores, color=colores)

    ax.set_title(
        titulo,
        fontsize=15,
        fontweight="bold"
    )

    ax.set_xticks([])
    ax.tick_params(axis="y", labelsize=11)

    max_val = max(valores) if max(valores) > 0 else 1

    for i, v in enumerate(valores):
        ax.text(
            v + (max_val * 0.01),
            i,
            f"{int(v):,}",
            va="center",
            fontsize=11,
            fontweight="bold"
        )

    ax.set_xlim(0, max_val * 1.18)
    ax.set_facecolor("none")
    fig.patch.set_alpha(0)

    plt.tight_layout()

    return fig

def crear_graficas(m):

    fig_piezas = grafica(
        "Piezas",
        [m["p"], m["e"], m["t"]],
        ["Emitido", "Entregado", "Tránsito"],
        ["#235B4E", "#9F2241", "#B38E5D"]
    )

    fig_montos = grafica(
        "Montos ($)",
        [m["m_p"], m["m_e"], m["m_t"]],
        ["Emitido", "Entregado", "Tránsito"],
        ["#235B4E", "#9F2241", "#B38E5D"]
    )

    fig_ordenes = grafica(
        "Órdenes",
        [m["o_total"], m["o_ent"], m["o_tran"]],
        ["Total", "Entregadas", "En tránsito"],
        ["#235B4E", "#9F2241", "#B38E5D"]
    )

    fig_claves = grafica(
        "Claves",
        [m["c_total"], m["c_ent"], m["c_tran"]],
        ["Total", "Entregadas", "En tránsito"],
        ["#235B4E", "#9F2241", "#B38E5D"]
    )

    return fig_piezas, fig_montos, fig_ordenes, fig_claves

# =========================
# RESUMEN GENERAL NACIONAL
# =========================
metricas_general = calcular_metricas(df)

st.markdown("## 🌎 Resumen general nacional")

g1, g2, g3 = st.columns(3)

with g1:
    st.metric(
        "🔵 Total emitido",
        fmt(metricas_general["p"])
    )

with g2:
    st.metric(
        "🟢 Total entregado",
        fmt(metricas_general["e"])
    )

with g3:
    st.metric(
        "🟡 Total en tránsito",
        fmt(metricas_general["t"])
    )

st.markdown("### 📋 Tabla nacional")

col_n1, col_n2, col_n3 = st.columns([1, 3, 1])

with col_n2:
    st.dataframe(
        metricas_general["tabla"],
        use_container_width=True,
        hide_index=True
    )

st.divider()

# =========================
# FILTRO
# =========================
estado_sel = st.selectbox(
    "📍 Estado",
    sorted(df[col_estado].dropna().unique())
)

df_f = df[df[col_estado] == estado_sel].copy()

metricas = calcular_metricas(df_f)

tabla_base = metricas["tabla"]

fig_piezas, fig_montos, fig_ordenes, fig_claves = crear_graficas(metricas)

# =========================
# TABLA STREAMLIT
# =========================
st.markdown(f"## 📋 Tabla operativa - {estado_sel}")

col1, col2, col3 = st.columns([1, 3, 1])

with col2:
    st.dataframe(
        tabla_base,
        use_container_width=True,
        hide_index=True
    )

# =========================
# VISUALIZACIÓN
# =========================
st.markdown(f"## 📊 Visualización - {estado_sel}")

c1, c2 = st.columns(2)
c3, c4 = st.columns(2)

with c1:
    st.pyplot(fig_piezas)

with c2:
    st.pyplot(fig_montos)

with c3:
    st.pyplot(fig_ordenes)

with c4:
    st.pyplot(fig_claves)

# =========================
# POWERPOINT
# =========================
MACHOTE = r"C:\Users\guillermo.ortega\OneDrive - IMSS-BIENESTAR\Escritorio\python\MACHOTE_PRESENTACIÓN.pptx"

def obtener_layout(prs):
    if len(prs.slide_layouts) > 5:
        return prs.slide_layouts[5]
    elif len(prs.slide_layouts) > 0:
        return prs.slide_layouts[0]
    else:
        raise Exception("El PowerPoint no tiene layouts disponibles.")

def limpiar_slide(slide):
    for shape in list(slide.shapes):
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception:
            pass

def agregar_contenido_slide(
    slide,
    tabla,
    fig_piezas,
    fig_montos,
    fig_ordenes,
    fig_claves,
    estado
):

    txBox = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.2),
        Inches(8),
        Inches(0.5)
    )

    tf = txBox.text_frame
    tf.text = f"Abasto - {str(estado).title()}"

    p_title = tf.paragraphs[0]
    p_title.font.size = Pt(30)
    p_title.font.bold = True

    fig_tab, ax = plt.subplots(figsize=(10, 2))
    ax.axis("off")

    table = ax.table(
        cellText=tabla.values,
        colLabels=tabla.columns,
        loc="center",
        cellLoc="center"
    )

    table.scale(1.4, 1.5)

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("black")

        if r == 0:
            cell.set_facecolor("#235B4E")
            cell.set_text_props(
                color="white",
                weight="bold",
                size=10
            )
        else:
            cell.set_facecolor((1, 1, 1, 0))
            cell.set_text_props(size=10)

    fig_tab.patch.set_alpha(0)

    img_tab = io.BytesIO()

    plt.savefig(
        img_tab,
        format="png",
        bbox_inches="tight",
        transparent=True
    )

    img_tab.seek(0)

    plt.close(fig_tab)

    slide.shapes.add_picture(
        img_tab,
        Inches(1.75),
        Inches(0.75),
        width=Inches(9.7)
    )

    def add(fig, x, y):

        img = io.BytesIO()

        fig.savefig(
            img,
            format="png",
            bbox_inches="tight",
            transparent=True
        )

        img.seek(0)

        slide.shapes.add_picture(
            img,
            Inches(x),
            Inches(y),
            width=Inches(4.9)
        )

        plt.close(fig)

    add(fig_piezas, 0.55, 2.05)
    add(fig_montos, 6.55, 2.05)
    add(fig_ordenes, 0.55, 4.25)
    add(fig_claves, 6.55, 4.25)

def exportar_ppt_estado(
    tabla_base,
    fig_piezas,
    fig_montos,
    fig_ordenes,
    fig_claves,
    estado_sel
):

    prs = Presentation(MACHOTE) if os.path.exists(MACHOTE) else Presentation()

    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(obtener_layout(prs))
    else:
        slide = prs.slides[0]

    limpiar_slide(slide)

    agregar_contenido_slide(
        slide,
        tabla_base,
        fig_piezas,
        fig_montos,
        fig_ordenes,
        fig_claves,
        estado_sel
    )

    salida = f"Dashboard_IMB_{estado_sel}.pptx"

    prs.save(salida)

    return salida

def exportar_ppt_todos_estados():

    prs = Presentation(MACHOTE) if os.path.exists(MACHOTE) else Presentation()

    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    estados = sorted(df[col_estado].dropna().unique())

    barra = st.progress(0)
    total = len(estados)

    for i, estado in enumerate(estados):

        df_estado = df[df[col_estado] == estado].copy()

        m = calcular_metricas(df_estado)

        fig_p, fig_m, fig_o, fig_c = crear_graficas(m)

        slide = prs.slides.add_slide(obtener_layout(prs))

        limpiar_slide(slide)

        agregar_contenido_slide(
            slide,
            m["tabla"],
            fig_p,
            fig_m,
            fig_o,
            fig_c,
            estado
        )

        barra.progress((i + 1) / total)

    salida = "Dashboard_Nacional_IMB_Todos_Estados.pptx"

    prs.save(salida)

    return salida

# =========================
# BOTONES
# =========================
st.markdown("## 📤 Exportación PPT")

col_exp1, col_exp2 = st.columns(2)

with col_exp1:

    if st.button("📊 Exportar estado seleccionado"):

        archivo_ppt = exportar_ppt_estado(
            tabla_base,
            fig_piezas,
            fig_montos,
            fig_ordenes,
            fig_claves,
            estado_sel
        )

        st.success("PPT del estado generado correctamente ✅")

        with open(archivo_ppt, "rb") as f:
            st.download_button(
                "⬇ Descargar PPT del estado",
                f,
                file_name=archivo_ppt
            )

with col_exp2:

    if st.button("🌎 Exportar TODOS los estados"):

        archivo_total = exportar_ppt_todos_estados()

        st.success("PowerPoint nacional generado correctamente ✅")

        with open(archivo_total, "rb") as f:
            st.download_button(
                "⬇ Descargar PowerPoint Nacional",
                f,
                file_name=archivo_total
            )
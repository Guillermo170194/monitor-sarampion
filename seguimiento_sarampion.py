import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# =========================
# CONFIG
# =========================
st.set_page_config(
    page_title="Monitor Nacional Sarampión",
    layout="wide"
)

# =========================
# ESTILO
# =========================
st.markdown("""
<style>
html, body, [class*="css"]{
    font-family:'Noto Sans', sans-serif;
}

.main-title{
    text-align:center;
    color:#7A1735;
    font-size:42px;
    font-weight:900;
}

.sub-title{
    text-align:center;
    color:#555;
    font-size:15px;
    margin-bottom:10px;
}

.fecha{
    text-align:center;
    color:#777;
    font-size:12px;
    margin-bottom:20px;
}

.card{
    padding:16px;
    border-radius:18px;
    background:white;
    box-shadow:0px 4px 12px rgba(0,0,0,0.10);
    text-align:center;
}

.azul{border-top:8px solid #2F5DA9;}
.amarillo{border-top:8px solid #B38E5D;}
.verde{border-top:8px solid #235B4E;}
.rojo{border-top:8px solid #9F2241;}

.kpi{
    font-size:30px;
    font-weight:900;
}

.label{
    font-size:13px;
    font-weight:700;
    color:#666;
}

.alerta{
    padding:8px 12px;
    border-radius:10px;
    background:#fff5f5;
    border-left:5px solid #9F2241;
    margin-bottom:6px;
    font-weight:600;
    font-size:13px;
}

.ok{
    padding:8px 12px;
    border-radius:10px;
    background:#f4fbf7;
    border-left:5px solid #235B4E;
    margin-bottom:6px;
    font-weight:600;
    font-size:13px;
}
</style>
""", unsafe_allow_html=True)

st.markdown(
    "<div class='main-title'>💉 Monitor Nacional de Sarampión</div>",
    unsafe_allow_html=True
)

st.markdown(
    "<div class='sub-title'>Seguimiento operativo nacional de claves estratégicas</div>",
    unsafe_allow_html=True
)

fecha_actual = datetime.now().strftime("%d/%m/%Y %H:%M hrs")

st.markdown(
    f"<div class='fecha'>Actualizado: {fecha_actual}</div>",
    unsafe_allow_html=True
)

# =========================
# RUTAS
# =========================

RUTA_BD = r"C:\Users\guillermo.ortega\IMSS-BIENESTAR\Emisión y Proveeduría - Consolidada 2026\Consolidada 2026 IMB.xlsb"

RUTA_INV = r"C:\Users\guillermo.ortega\OneDrive - IMSS-BIENESTAR\Escritorio\python\INV 08 DE MAYO 2026 PREELIMINAR (1).xlsx"

RUTA_CPM = r"C:\Users\guillermo.ortega\OneDrive - IMSS-BIENESTAR\Escritorio\python\CPM CLUES Demanda IB Nov 2025 20251127 ACTUAL (4).xlsb"

# =========================
# CLAVES
# =========================
CLAVES = [
    "060.904.0100",
    "060.040.3711",
    "060.550.2657",
    "060.550.0438",
    "020.000.3820.00",
    "060.550.2707",
    "060.550.0636",
    "020.000.3800.00",
    "060.066.0039",
    "060.125.2505",
    "060.125.3958",
    "060.218.0085"
]

# =========================
# FUNCIONES
# =========================
def fmt(x):
    return f"{int(float(x)):,}"

def fmt_dec(x):
    return f"{float(x):,.2f}"

def normaliza(txt):
    if pd.isna(txt):
        return ""

    txt = str(txt).strip().upper()

    for a, b in {
        "Á": "A",
        "É": "E",
        "Í": "I",
        "Ó": "O",
        "Ú": "U"
    }.items():
        txt = txt.replace(a, b)

    if txt == "MICHOACAN":
        txt = "MICHOACAN DE OCAMPO"

    return txt

def limpiar_modelo(txt):
    txt = normaliza(txt)

    if "NO" in txt and "CONCURRENTE" in txt:
        return "NO CONCURRENTE"

    if "IMSS" in txt or "BIENESTAR" in txt or "CONCURRENTE" in txt:
        return "CONCURRENTE"

    return "SIN CLASIFICAR"

def resumir_descripcion(txt, max_palabras=8):
    if pd.isna(txt):
        return ""

    txt = str(txt).strip().upper()

    if txt in ["", "NAN", "NONE"]:
        return ""

    return " ".join(txt.split()[:max_palabras])

def clasificar_abasto(nivel, inventario, cpm):
    if inventario == 0:
        return "🔴 AGOTADO"

    if cpm == 0 and inventario > 0:
        return "🟢 ÓPTIMO"

    if nivel > 0 and nivel < 1:
        return "🟠 PRÓX AGOTARSE"

    if nivel >= 1 and nivel < 1.5:
        return "🟡 BAJO"

    if nivel >= 1.5 and nivel <= 5:
        return "🟢 ÓPTIMO"

    return "🔵 SOBRE DISPOSICIÓN"

def clasificar_proveedor(cumplimiento):
    if cumplimiento == 0:
        return "🔴 SIN ENTREGA"

    if cumplimiento < 0.5:
        return "🟠 BAJO CUMPLIMIENTO"

    if cumplimiento < 0.8:
        return "🟡 CUMPLIMIENTO PARCIAL"

    if cumplimiento <= 1:
        return "🟢 BUEN CUMPLIMIENTO"

    return "🔵 ENTREGA MAYOR A EMISIÓN"

def formatear_tabla(df):
    df_fmt = df.copy()

    cols_enteros = [
        "PIEZAS EMITIDAS",
        "PIEZAS ENTREGADAS",
        "PIEZAS EN TRÁNSITO",
        "PIEZAS EN INVENTARIO",
        "CPM",
        "TOTAL",
        "CLAVES",
        "CLAVES AGOTADAS"
    ]

    for c in cols_enteros:
        if c in df_fmt.columns:
            df_fmt[c] = df_fmt[c].apply(lambda x: f"{int(float(x)):,}")

    if "NIVEL DE ABASTO" in df_fmt.columns:
        df_fmt["NIVEL DE ABASTO"] = df_fmt["NIVEL DE ABASTO"].apply(
            lambda x: f"{float(x):,.2f}"
        )

    if "NIVEL DE CUMPLIMIENTO" in df_fmt.columns:
        df_fmt["NIVEL DE CUMPLIMIENTO"] = df_fmt["NIVEL DE CUMPLIMIENTO"].apply(
            lambda x: f"{float(x):,.2f}"
        )

    if "% CUMPLIMIENTO" in df_fmt.columns:
        df_fmt["% CUMPLIMIENTO"] = df_fmt["% CUMPLIMIENTO"].apply(
            lambda x: f"{float(x):,.2f}%"
        )

    return df_fmt

# =========================
# CARGA
# =========================
@st.cache_data(show_spinner=True)
def cargar():

    bd = pd.read_excel(
        RUTA_BD,
        sheet_name="BD",
        engine="pyxlsb",
        usecols=[3, 5, 8, 9, 15, 17, 18, 27, 55]
    )

    bd.columns = bd.columns.str.strip()

    COL_ESTADO = bd.columns[0]
    COL_MODELO = bd.columns[1]
    COL_PROVEEDOR = bd.columns[4]
    COL_CLAVE = bd.columns[5]
    COL_DESCRIPCION = bd.columns[6]
    COL_EMITIDAS = bd.columns[7]
    COL_ENTREGADAS = bd.columns[8]

    bd[COL_CLAVE] = bd[COL_CLAVE].astype(str).str.strip()
    bd = bd[bd[COL_CLAVE].isin(CLAVES)].copy()

    bd["ENTIDAD"] = bd[COL_ESTADO].apply(normaliza)
    bd["MODELO OPERATIVO"] = bd[COL_MODELO].apply(limpiar_modelo)
    bd["PROVEEDOR"] = bd[COL_PROVEEDOR].astype(str).str.strip().str.upper()
    bd["CLAVE"] = bd[COL_CLAVE].astype(str).str.strip()
    bd["DESCRIPCIÓN"] = bd[COL_DESCRIPCION].apply(resumir_descripcion)

    bd["PIEZAS EMITIDAS"] = pd.to_numeric(
        bd[COL_EMITIDAS],
        errors="coerce"
    ).fillna(0)

    bd["PIEZAS ENTREGADAS"] = pd.to_numeric(
        bd[COL_ENTREGADAS],
        errors="coerce"
    ).fillna(0)

    bd_group = (
        bd.groupby(
            [
                "ENTIDAD",
                "MODELO OPERATIVO",
                "PROVEEDOR",
                "CLAVE",
                "DESCRIPCIÓN"
            ],
            dropna=False
        )
        .agg({
            "PIEZAS EMITIDAS": "sum",
            "PIEZAS ENTREGADAS": "sum"
        })
        .reset_index()
    )

    inv = pd.read_excel(
        RUTA_INV,
        sheet_name="Resultado consulta"
    )

    inv.columns = inv.columns.str.strip()

    INV_ESTADO = inv.columns[0]
    INV_CLAVE = inv.columns[3]
    INV_EXIST = inv.columns[5]

    inv[INV_CLAVE] = inv[INV_CLAVE].astype(str).str.strip()
    inv = inv[inv[INV_CLAVE].isin(CLAVES)].copy()

    inv["ENTIDAD"] = inv[INV_ESTADO].apply(normaliza)
    inv["CLAVE"] = inv[INV_CLAVE].astype(str).str.strip()

    inv["PIEZAS EN INVENTARIO"] = pd.to_numeric(
        inv[INV_EXIST],
        errors="coerce"
    ).fillna(0)

    inv_group = (
        inv.groupby(
            ["ENTIDAD", "CLAVE"],
            dropna=False
        )["PIEZAS EN INVENTARIO"]
        .sum()
        .reset_index()
    )

    cpm = pd.read_excel(
        RUTA_CPM,
        sheet_name="cpm",
        engine="pyxlsb",
        usecols=[0, 1, 3]
    )

    cpm.columns = cpm.columns.str.strip()

    CPM_ESTADO = cpm.columns[0]
    CPM_CLAVE = cpm.columns[1]
    CPM_VALOR = cpm.columns[2]

    cpm["ENTIDAD"] = cpm[CPM_ESTADO].apply(normaliza)
    cpm["CLAVE"] = cpm[CPM_CLAVE].astype(str).str.strip()

    cpm["CPM"] = pd.to_numeric(
        cpm[CPM_VALOR],
        errors="coerce"
    ).fillna(0)

    cpm = cpm[cpm["CLAVE"].isin(CLAVES)].copy()

    cpm_group = (
        cpm.groupby(
            ["ENTIDAD", "CLAVE"],
            dropna=False
        )["CPM"]
        .sum()
        .reset_index()
    )

    df = bd_group.merge(
        inv_group,
        how="left",
        on=["ENTIDAD", "CLAVE"]
    )

    df = df.merge(
        cpm_group,
        how="left",
        on=["ENTIDAD", "CLAVE"]
    )

    df["PIEZAS EN INVENTARIO"] = df["PIEZAS EN INVENTARIO"].fillna(0)
    df["CPM"] = df["CPM"].fillna(0)

    df["PIEZAS EN TRÁNSITO"] = (
        df["PIEZAS EMITIDAS"] -
        df["PIEZAS ENTREGADAS"]
    )

    df.loc[
        df["PIEZAS EN TRÁNSITO"] < 0,
        "PIEZAS EN TRÁNSITO"
    ] = 0

    df["NIVEL DE ABASTO"] = 0.0

    mask = df["CPM"] > 0

    df.loc[mask, "NIVEL DE ABASTO"] = (
        df.loc[mask, "PIEZAS EN INVENTARIO"] /
        df.loc[mask, "CPM"]
    ).round(2)

    df["CLASIFICACION ABASTO"] = df.apply(
        lambda row: clasificar_abasto(
            row["NIVEL DE ABASTO"],
            row["PIEZAS EN INVENTARIO"],
            row["CPM"]
        ),
        axis=1
    )

    df["CLAVE FILTRO"] = (
        df["CLAVE"].astype(str) + " - " + df["DESCRIPCIÓN"].astype(str)
    )

    return df

# =========================
# CARGAR
# =========================
try:
    df = cargar()

except PermissionError:
    st.error("Cierra los archivos Excel y vuelve a ejecutar.")
    st.stop()

except Exception as e:
    st.error(f"Error: {e}")
    st.stop()

# =========================
# FILTROS
# =========================
st.sidebar.title("🔎 Filtros")

modelo = st.sidebar.multiselect(
    "Modelo operativo",
    sorted(df["MODELO OPERATIVO"].dropna().unique())
)

df_op = df.copy()

if modelo:
    df_op = df_op[df_op["MODELO OPERATIVO"].isin(modelo)]

entidad = st.sidebar.multiselect(
    "Entidad",
    sorted(df_op["ENTIDAD"].dropna().unique())
)

if entidad:
    df_op = df_op[df_op["ENTIDAD"].isin(entidad)]

proveedor = st.sidebar.multiselect(
    "Proveedor",
    sorted(df_op["PROVEEDOR"].dropna().unique())
)

if proveedor:
    df_op = df_op[df_op["PROVEEDOR"].isin(proveedor)]

clave_filtro = st.sidebar.multiselect(
    "Clave",
    sorted(df_op["CLAVE FILTRO"].dropna().unique())
)

clave = []

if clave_filtro:
    clave = [
        x.split(" - ")[0]
        for x in clave_filtro
    ]

    df_op = df_op[df_op["CLAVE"].isin(clave)]

clasif = st.sidebar.multiselect(
    "Clasificación abasto",
    sorted(df_op["CLASIFICACION ABASTO"].dropna().unique())
)

df_f = df.copy()

if modelo:
    df_f = df_f[df_f["MODELO OPERATIVO"].isin(modelo)]

if entidad:
    df_f = df_f[df_f["ENTIDAD"].isin(entidad)]

if proveedor:
    df_f = df_f[df_f["PROVEEDOR"].isin(proveedor)]

if clave:
    df_f = df_f[df_f["CLAVE"].isin(clave)]

if clasif:
    df_f = df_f[df_f["CLASIFICACION ABASTO"].isin(clasif)]

# =========================
# KPI
# =========================
emitido = df_f["PIEZAS EMITIDAS"].sum()
entregado = df_f["PIEZAS ENTREGADAS"].sum()
transito = df_f["PIEZAS EN TRÁNSITO"].sum()
inventario = df_f["PIEZAS EN INVENTARIO"].sum()

agotadas = (
    df_f[df_f["CLASIFICACION ABASTO"] == "🔴 AGOTADO"]["CLAVE"]
    .nunique()
)

k1, k2, k3, k4, k5 = st.columns(5)

k1.markdown(f"""
<div class='card azul'>
<div class='label'>🔵 Piezas emitidas</div>
<div class='kpi'>{fmt(emitido)}</div>
</div>
""", unsafe_allow_html=True)

k2.markdown(f"""
<div class='card verde'>
<div class='label'>🟢 Piezas entregadas</div>
<div class='kpi'>{fmt(entregado)}</div>
</div>
""", unsafe_allow_html=True)

k3.markdown(f"""
<div class='card amarillo'>
<div class='label'>🟡 Piezas en tránsito</div>
<div class='kpi'>{fmt(transito)}</div>
</div>
""", unsafe_allow_html=True)

k4.markdown(f"""
<div class='card verde'>
<div class='label'>🟢 Piezas en inventario</div>
<div class='kpi'>{fmt(inventario)}</div>
</div>
""", unsafe_allow_html=True)

k5.markdown(f"""
<div class='card rojo'>
<div class='label'>🚨 Claves agotadas</div>
<div class='kpi'>{fmt(agotadas)}</div>
</div>
""", unsafe_allow_html=True)

st.divider()

# =========================
# ALERTAS
# =========================
st.markdown("### 🚨 Alertas prioritarias")

alertas_df = df_f[
    df_f["CLASIFICACION ABASTO"] == "🔴 AGOTADO"
]

if alertas_df.empty:

    st.markdown(
        """
        <div class='ok'>
        Sin claves agotadas.
        </div>
        """,
        unsafe_allow_html=True
    )

else:

    top_alertas = (
        alertas_df.groupby("ENTIDAD")["CLAVE"]
        .nunique()
        .reset_index(name="TOTAL")
        .sort_values("TOTAL", ascending=False)
    )

    col1, col2, col3 = st.columns(3)
    columnas = [col1, col2, col3]

    for i, (_, row) in enumerate(top_alertas.iterrows()):

        with columnas[i % 3]:

            st.markdown(
                f"""
                <div class='alerta'>
                <b>{row['ENTIDAD']}</b><br>
                {fmt(row['TOTAL'])} claves agotadas
                </div>
                """,
                unsafe_allow_html=True
            )

# =========================
# TABS
# =========================
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "🗺️ Entidades",
    "💉 Claves",
    "📊 Gráficas",
    "📋 Detalle",
    "🏭 Proveedores",
    "📥 PowerPoint"
])

# =========================
# TAB ENTIDADES
# =========================
with tab1:

    st.markdown("## 🗺️ Análisis por entidad")

    resumen = (
        df_f.groupby(
            ["ENTIDAD", "MODELO OPERATIVO"]
        )
        .agg({
            "PIEZAS EMITIDAS": "sum",
            "PIEZAS ENTREGADAS": "sum",
            "PIEZAS EN TRÁNSITO": "sum",
            "PIEZAS EN INVENTARIO": "sum",
            "CPM": "sum"
        })
        .reset_index()
    )

    resumen["NIVEL DE ABASTO"] = 0.0

    mask = resumen["CPM"] > 0

    resumen.loc[mask, "NIVEL DE ABASTO"] = (
        resumen.loc[mask, "PIEZAS EN INVENTARIO"] /
        resumen.loc[mask, "CPM"]
    ).round(2)

    resumen["CLASIFICACION ABASTO"] = resumen.apply(
        lambda row: clasificar_abasto(
            row["NIVEL DE ABASTO"],
            row["PIEZAS EN INVENTARIO"],
            row["CPM"]
        ),
        axis=1
    )

    agot = (
        df_f[df_f["CLASIFICACION ABASTO"] == "🔴 AGOTADO"]
        .groupby("ENTIDAD")["CLAVE"]
        .nunique()
        .reset_index(name="CLAVES AGOTADAS")
    )

    resumen = resumen.merge(
        agot,
        how="left",
        on="ENTIDAD"
    )

    resumen["CLAVES AGOTADAS"] = (
        resumen["CLAVES AGOTADAS"]
        .fillna(0)
        .astype(int)
    )

    resumen = resumen.sort_values(
        ["CLAVES AGOTADAS", "NIVEL DE ABASTO"],
        ascending=[False, True]
    )

    st.dataframe(
        formatear_tabla(resumen),
        use_container_width=True,
        hide_index=True
    )

    if len(entidad) == 1:

        st.markdown("### 💉 Tabla por clave")

        tabla_estado = (
            df_f.groupby(
                ["CLAVE", "DESCRIPCIÓN", "CLASIFICACION ABASTO"]
            )
            .agg({
                "PIEZAS EMITIDAS": "sum",
                "PIEZAS ENTREGADAS": "sum",
                "PIEZAS EN TRÁNSITO": "sum",
                "PIEZAS EN INVENTARIO": "sum",
                "CPM": "sum"
            })
            .reset_index()
        )

        tabla_estado["NIVEL DE ABASTO"] = 0.0

        mask_estado = tabla_estado["CPM"] > 0

        tabla_estado.loc[mask_estado, "NIVEL DE ABASTO"] = (
            tabla_estado.loc[mask_estado, "PIEZAS EN INVENTARIO"] /
            tabla_estado.loc[mask_estado, "CPM"]
        ).round(2)

        tabla_estado["CLASIFICACION ABASTO"] = tabla_estado.apply(
            lambda row: clasificar_abasto(
                row["NIVEL DE ABASTO"],
                row["PIEZAS EN INVENTARIO"],
                row["CPM"]
            ),
            axis=1
        )

        tabla_estado = tabla_estado.sort_values(
            ["NIVEL DE ABASTO"],
            ascending=True
        )

        st.dataframe(
            formatear_tabla(tabla_estado),
            use_container_width=True,
            hide_index=True
        )

# =========================
# TAB CLAVES
# =========================
with tab2:

    st.markdown("## 💉 Estatus por clave")

    tabla = (
        df_f.groupby(
            ["CLAVE", "DESCRIPCIÓN", "CLASIFICACION ABASTO"]
        )
        .agg({
            "PIEZAS EMITIDAS": "sum",
            "PIEZAS ENTREGADAS": "sum",
            "PIEZAS EN TRÁNSITO": "sum",
            "PIEZAS EN INVENTARIO": "sum",
            "CPM": "sum"
        })
        .reset_index()
    )

    tabla["NIVEL DE ABASTO"] = 0.0

    mask = tabla["CPM"] > 0

    tabla.loc[mask, "NIVEL DE ABASTO"] = (
        tabla.loc[mask, "PIEZAS EN INVENTARIO"] /
        tabla.loc[mask, "CPM"]
    ).round(2)

    tabla["CLASIFICACION ABASTO"] = tabla.apply(
        lambda row: clasificar_abasto(
            row["NIVEL DE ABASTO"],
            row["PIEZAS EN INVENTARIO"],
            row["CPM"]
        ),
        axis=1
    )

    tabla = tabla.sort_values(
        ["NIVEL DE ABASTO"],
        ascending=True
    )

    st.dataframe(
        formatear_tabla(tabla),
        use_container_width=True,
        hide_index=True
    )

# =========================
# TAB GRAFICAS
# =========================
with tab3:

    st.markdown("## 📊 Gráficas por clave y piezas")

    graf_clave = (
        df_f.groupby(["CLAVE", "DESCRIPCIÓN"])
        .agg({
            "PIEZAS EMITIDAS": "sum",
            "PIEZAS ENTREGADAS": "sum",
            "PIEZAS EN TRÁNSITO": "sum",
            "PIEZAS EN INVENTARIO": "sum",
            "CPM": "sum"
        })
        .reset_index()
    )

    graf_clave["CLAVE_DESC"] = (
        graf_clave["CLAVE"] + " - " + graf_clave["DESCRIPCIÓN"]
    )

    fig1, ax1 = plt.subplots(figsize=(12, 6))

    ax1.barh(
        graf_clave["CLAVE_DESC"],
        graf_clave["PIEZAS EMITIDAS"],
        label="Piezas emitidas"
    )

    ax1.barh(
        graf_clave["CLAVE_DESC"],
        graf_clave["PIEZAS ENTREGADAS"],
        label="Piezas entregadas"
    )

    ax1.barh(
        graf_clave["CLAVE_DESC"],
        graf_clave["PIEZAS EN TRÁNSITO"],
        label="Piezas en tránsito"
    )

    ax1.barh(
        graf_clave["CLAVE_DESC"],
        graf_clave["PIEZAS EN INVENTARIO"],
        label="Piezas en inventario"
    )

    ax1.set_title(
        "Piezas emitidas, entregadas, tránsito e inventario por clave",
        fontweight="bold"
    )

    ax1.set_xticks([])
    ax1.legend()

    plt.tight_layout()
    st.pyplot(fig1)

    st.markdown("### 📈 CPM vs inventario por clave")

    fig2, ax2 = plt.subplots(figsize=(12, 6))

    ax2.barh(
        graf_clave["CLAVE_DESC"],
        graf_clave["CPM"],
        label="CPM"
    )

    ax2.barh(
        graf_clave["CLAVE_DESC"],
        graf_clave["PIEZAS EN INVENTARIO"],
        label="Piezas en inventario"
    )

    ax2.set_title(
        "CPM vs piezas en inventario por clave",
        fontweight="bold"
    )

    ax2.set_xticks([])
    ax2.legend()

    plt.tight_layout()
    st.pyplot(fig2)

# =========================
# TAB DETALLE
# =========================
with tab4:

    st.markdown("## 📋 Detalle operativo")

    columnas_detalle = [
        "ENTIDAD",
        "MODELO OPERATIVO",
        "PROVEEDOR",
        "CLAVE",
        "DESCRIPCIÓN",
        "PIEZAS EMITIDAS",
        "PIEZAS ENTREGADAS",
        "PIEZAS EN TRÁNSITO",
        "PIEZAS EN INVENTARIO",
        "CPM",
        "NIVEL DE ABASTO",
        "CLASIFICACION ABASTO"
    ]

    st.dataframe(
        formatear_tabla(df_f[columnas_detalle]),
        use_container_width=True,
        hide_index=True
    )

    excel = io.BytesIO()

    with pd.ExcelWriter(
        excel,
        engine="openpyxl"
    ) as writer:
        df_f[columnas_detalle].to_excel(
            writer,
            index=False,
            sheet_name="DETALLE"
        )

    excel.seek(0)

    st.download_button(
        "⬇ Descargar Excel",
        excel,
        file_name="monitor_sarampion.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

# =========================
# TAB PROVEEDORES
# =========================
with tab5:

    st.markdown("## 🏭 Nivel de abasto por proveedor")

    st.write(
        "Este apartado mide el cumplimiento del proveedor con base en piezas entregadas contra piezas emitidas."
    )

    if st.button("📊 Calcular abasto de proveedores"):

        proveedores = (
            df_f.groupby("PROVEEDOR")
            .agg({
                "PIEZAS EMITIDAS": "sum",
                "PIEZAS ENTREGADAS": "sum",
                "PIEZAS EN TRÁNSITO": "sum",
                "PIEZAS EN INVENTARIO": "sum",
                "CLAVE": "nunique"
            })
            .reset_index()
            .rename(columns={"CLAVE": "CLAVES"})
        )

        proveedores["NIVEL DE CUMPLIMIENTO"] = 0.0

        mask_prov = proveedores["PIEZAS EMITIDAS"] > 0

        proveedores.loc[mask_prov, "NIVEL DE CUMPLIMIENTO"] = (
            proveedores.loc[mask_prov, "PIEZAS ENTREGADAS"] /
            proveedores.loc[mask_prov, "PIEZAS EMITIDAS"]
        ).round(2)

        proveedores["% CUMPLIMIENTO"] = (
            proveedores["NIVEL DE CUMPLIMIENTO"] * 100
        ).round(2)

        proveedores["CLASIFICACIÓN PROVEEDOR"] = proveedores[
            "NIVEL DE CUMPLIMIENTO"
        ].apply(clasificar_proveedor)

        proveedores = proveedores.sort_values(
            "% CUMPLIMIENTO",
            ascending=True
        )

        st.session_state["proveedores_abasto"] = proveedores

    if "proveedores_abasto" in st.session_state:

        proveedores = st.session_state["proveedores_abasto"]

        st.dataframe(
            formatear_tabla(proveedores),
            use_container_width=True,
            hide_index=True
        )

        excel_prov = io.BytesIO()

        with pd.ExcelWriter(
            excel_prov,
            engine="openpyxl"
        ) as writer:
            proveedores.to_excel(
                writer,
                index=False,
                sheet_name="PROVEEDORES"
            )

        excel_prov.seek(0)

        st.download_button(
            "⬇ Descargar análisis de proveedores",
            excel_prov,
            file_name="nivel_abasto_proveedores.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

        st.markdown("### 📊 Cumplimiento por proveedor")

        graf_prov = proveedores.head(20).copy()

        figp, axp = plt.subplots(figsize=(12, 7))

        axp.barh(
            graf_prov["PROVEEDOR"],
            graf_prov["% CUMPLIMIENTO"]
        )

        axp.set_title(
            "Nivel de cumplimiento por proveedor",
            fontweight="bold"
        )

        axp.set_xlabel("% cumplimiento")
        plt.tight_layout()

        st.pyplot(figp)

# =========================
# POWERPOINT
# =========================
import os

def agregar_texto(slide, texto, x, y, w, h, size=18, bold=False, color=(0, 0, 0)):

    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()

    run.text = texto
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Noto Sans"
    run.font.color.rgb = RGBColor(*color)

    return box


def agregar_tabla(
    slide,
    df_tabla,
    x,
    y,
    w,
    h,
    font_header=8,
    font_body=7
):

    rows = len(df_tabla) + 1
    cols = len(df_tabla.columns)

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    table = table_shape.table

    # encabezados
    for c, col in enumerate(df_tabla.columns):

        cell = table.cell(0, c)

        cell.text = str(col)

        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(35, 91, 78)

        for p in cell.text_frame.paragraphs:
            for run in p.runs:

                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(font_header)
                run.font.name = "Noto Sans"

    # datos
    for r in range(len(df_tabla)):

        for c, col in enumerate(df_tabla.columns):

            cell = table.cell(r + 1, c)

            cell.text = str(df_tabla.iloc[r, c])

            for p in cell.text_frame.paragraphs:
                for run in p.runs:

                    run.font.size = Pt(font_body)
                    run.font.name = "Noto Sans"

    return table_shape


def generar_ppt(df_base):

    RUTA_MACHOTE = r"C:\Users\guillermo.ortega\OneDrive - IMSS-BIENESTAR\Escritorio\python\MACHOTE_PRESENTACIÓN.pptx"

    if os.path.exists(RUTA_MACHOTE):

        prs = Presentation(RUTA_MACHOTE)

    else:

        prs = Presentation()

        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[-1]

    entidades = sorted(
        df_base["ENTIDAD"]
        .dropna()
        .unique()
    )

    for ent in entidades:

        dfe = df_base[
            df_base["ENTIDAD"] == ent
        ].copy()

        slide = prs.slides.add_slide(blank)

        # limpiar shapes del layout
        for shape in list(slide.shapes):

            sp = shape._element
            sp.getparent().remove(sp)

        modelo_ent = dfe["MODELO OPERATIVO"].iloc[0]

        # =========================
        # TITULO
        # =========================
        agregar_texto(
            slide,
            f"Seguimiento Sarampión - {ent.title()}",
            0.4,
            0.2,
            12,
            0.4,
            size=24,
            bold=True,
            color=(122, 23, 53)
        )

        agregar_texto(
            slide,
            f"Modelo operativo: {modelo_ent}",
            0.4,
            0.62,
            12,
            0.3,
            size=13,
            bold=True,
            color=(35, 91, 78)
        )

        # =========================
        # RESUMEN
        # =========================
        resumen = pd.DataFrame({

            "Indicador": [
                "Emitidas",
                "Entregadas",
                "Tránsito",
                "Inventario",
                "CPM",
                "Abasto"
            ],

            "Valor": [
                fmt(dfe["PIEZAS EMITIDAS"].sum()),
                fmt(dfe["PIEZAS ENTREGADAS"].sum()),
                fmt(dfe["PIEZAS EN TRÁNSITO"].sum()),
                fmt(dfe["PIEZAS EN INVENTARIO"].sum()),
                fmt(dfe["CPM"].sum()),
                fmt_dec(dfe["NIVEL DE ABASTO"].mean())
            ]
        })

        agregar_tabla(
            slide,
            resumen,
            0.4,
            1.0,
            3.2,
            2.2,
            font_header=10,
            font_body=9
        )

        # =========================
        # TABLA CLAVES
        # =========================
        tabla = dfe[[
            "CLAVE",
            "DESCRIPCIÓN",
            "PIEZAS ENTREGADAS",
            "PIEZAS EN INVENTARIO",
            "CPM",
            "NIVEL DE ABASTO",
            "CLASIFICACION ABASTO"
        ]].copy()

        tabla = tabla.rename(columns={
            "PIEZAS ENTREGADAS": "ENTREGADAS",
            "PIEZAS EN INVENTARIO": "INVENTARIO",
            "NIVEL DE ABASTO": "ABASTO",
            "CLASIFICACION ABASTO": "ESTATUS"
        })

        # resumir descripcion
        tabla["DESCRIPCIÓN"] = tabla["DESCRIPCIÓN"].astype(str).apply(
            lambda x: " ".join(x.split()[:4])
        )

        # quitar emojis
        tabla["ESTATUS"] = tabla["ESTATUS"].astype(str)

        for emo in ["🔴 ", "🟠 ", "🟡 ", "🟢 ", "🔵 "]:

            tabla["ESTATUS"] = tabla["ESTATUS"].str.replace(
                emo,
                "",
                regex=False
            )

        # formatos
        for col in [
            "ENTREGADAS",
            "INVENTARIO",
            "CPM"
        ]:

            tabla[col] = tabla[col].apply(
                lambda x: f"{int(float(x)):,}"
            )

        tabla["ABASTO"] = tabla["ABASTO"].apply(
            lambda x: f"{float(x):,.2f}"
        )

        tabla = tabla.sort_values(
            "ESTATUS"
        )

        agregar_texto(
            slide,
            "Estatus resumido por clave",
            3.9,
            1.0,
            8,
            0.3,
            size=15,
            bold=True,
            color=(35, 91, 78)
        )

        agregar_tabla(
            slide,
            tabla,
            3.9,
            1.35,
            8.9,
            4.9,
            font_header=8,
            font_body=7
        )

    salida = io.BytesIO()

    prs.save(salida)

    salida.seek(0)

    return salida


# =========================
# TAB POWERPOINT
# =========================
with tab6:

    st.markdown("## 📥 PowerPoint ejecutivo por entidad")

    st.write(
        "Genera un reporte ejecutivo resumido por entidad."
    )

    if st.button("📊 Generar PowerPoint"):

        with st.spinner("Generando PowerPoint..."):

            st.session_state["ppt_sarampion"] = generar_ppt(df_f)

        st.success("PowerPoint generado correctamente ✅")

    if "ppt_sarampion" in st.session_state:

        st.download_button(
            "⬇ Descargar PowerPoint ejecutivo",
            data=st.session_state["ppt_sarampion"],
            file_name="Reporte_Ejecutivo_Sarampion_por_Entidad.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )